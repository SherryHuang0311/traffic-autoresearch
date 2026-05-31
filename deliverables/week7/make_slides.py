"""
Generate final_presentation.pptx — 8 slides, ~5 min talk.
Run: python3 deliverables/week7/make_slides.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

OUTPUT = "deliverables/week7/final_presentation.pptx"

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1a, 0x3a, 0x5c)
BLUE    = RGBColor(0x2e, 0x6d, 0xa4)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LGREY   = RGBColor(0xf2, 0xf4, 0xf8)
DGREY   = RGBColor(0x44, 0x44, 0x44)
GREEN   = RGBColor(0x1d, 0x7a, 0x3a)
LGREEN  = RGBColor(0xd4, 0xed, 0xda)
RED     = RGBColor(0xa3, 0x16, 0x21)
LRED    = RGBColor(0xf8, 0xd7, 0xda)
AMBER   = RGBColor(0x85, 0x64, 0x04)
LAMBER  = RGBColor(0xff, 0xf3, 0xcd)
LBLUE   = RGBColor(0xd0, 0xe4, 0xf5)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rgb_hex(c: RGBColor):
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def bg(slide, color: RGBColor):
    """Set slide background colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          size=20, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def add_text_to_shape(shape, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def header_bar(slide, title, subtitle=None):
    """Dark navy bar at top."""
    bar = rect(slide, 0, 0, W, Inches(1.2), fill=NAVY)
    txbox(slide, title, Inches(0.35), Inches(0.15), Inches(12.0), Inches(0.65),
          size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle, Inches(0.35), Inches(0.78), Inches(12.0), Inches(0.35),
              size=13, bold=False, color=LBLUE, align=PP_ALIGN.LEFT)

def slide_num(slide, n, total=8):
    txbox(slide, f"{n} / {total}", Inches(12.5), Inches(7.15), Inches(0.8), Inches(0.3),
          size=10, color=DGREY, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, l, t, w, h, size=17, color=DGREY, gap=0.05):
    """Simple bulleted text box."""
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.text_frame.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.text_frame.paragraphs[0]
            first = False
        else:
            p = tf.text_frame.add_paragraph()
        p.space_before = Pt(gap * 72)
        run = p.add_run()
        bold_part = ""
        rest = item
        if item.startswith("**") and "**" in item[2:]:
            end = item.index("**", 2)
            bold_part = item[2:end]
            rest = item[end+2:]
        if bold_part:
            r1 = p.add_run()
            r1.text = "• " + bold_part
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            run.text = rest
        else:
            run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def mini_table(slide, rows, col_widths, l, t, row_h=Inches(0.38),
               header_fill=NAVY, alt_fill=LGREY, text_size=14):
    """Draw a simple table as coloured rectangles + text."""
    for ri, row in enumerate(rows):
        x = l
        fill = header_fill if ri == 0 else (alt_fill if ri % 2 == 0 else WHITE)
        txt_color = WHITE if ri == 0 else DGREY
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            r = rect(slide, x, t + ri*row_h, cw, row_h, fill=fill)
            txbox(slide, str(cell), x + Inches(0.08), t + ri*row_h + Inches(0.06),
                  cw - Inches(0.12), row_h - Inches(0.1),
                  size=text_size, bold=(ri == 0), color=txt_color,
                  align=PP_ALIGN.CENTER)
            x += cw

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title + Claim
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
rect(s, 0, Inches(2.5), W, Inches(2.8), fill=BLUE)

txbox(s, "AutoResearch for Traffic Congestion Prediction",
      Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.9),
      size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "Segment-relative features break a labeling-driven F1 ceiling",
      Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
      size=20, bold=False, color=LBLUE, align=PP_ALIGN.CENTER, italic=True)

txbox(s, "STAT 390 Capstone  •  Sherry Huang  •  Northwestern University",
      Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.45),
      size=16, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "One-sentence claim:",
      Inches(1.2), Inches(3.3), Inches(10.9), Inches(0.4),
      size=14, bold=True, color=LAMBER, align=PP_ALIGN.LEFT)

txbox(s, "Adding 4 road-specific speed features to a LightGBM model raises 30-min-ahead "
         "congestion F1 from 0.560 → 0.681 on validation (+21.6%) and 0.618 on held-out test, "
         "breaking a ceiling that 18 hyperparameter experiments could not cross.",
      Inches(1.2), Inches(3.7), Inches(10.9), Inches(1.3),
      size=17, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

txbox(s, "30 experiments  •  0 crashes  •  Test set opened exactly once",
      Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
      size=13, color=LBLUE, align=PP_ALIGN.CENTER)
slide_num(s, 1)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Topic → AutoResearch Contract
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LGREY)
header_bar(s, "From Topic to AutoResearch Contract",
           "How the research question was formalized")
slide_num(s, 2)

# Left: question
rect(s, Inches(0.3), Inches(1.35), Inches(5.9), Inches(5.6), fill=WHITE)
txbox(s, "The Research Question", Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.45),
      size=16, bold=True, color=NAVY)
items_q = [
    "Can speed history alone predict congestion 30 min ahead?",
    "No real-time feeds, no external map data",
    "Binary classification per road segment",
    "Chicago Traffic Tracker: 786K rows, Sept 2023",
    "Metric: validation binary F1 (positive = congested)",
]
bullet_box(s, items_q, Inches(0.5), Inches(2.05), Inches(5.6), Inches(3.5), size=16, color=DGREY)

# Right: contract terms
rect(s, Inches(6.5), Inches(1.35), Inches(6.5), Inches(5.6), fill=NAVY)
txbox(s, "AutoResearch Contract", Inches(6.7), Inches(1.5), Inches(6.1), Inches(0.45),
      size=16, bold=True, color=WHITE)
items_c = [
    "Modify ONLY src/model.py",
    "src/run.py FROZEN — split, label, metric immutable",
    "Keep if val F1 improves, revert if not",
    "Git commit every improvement, git checkout every discard",
    "Test set locked — opened exactly once at the end",
    "All 30 runs logged to experiments/results.csv",
]
bullet_box(s, items_c, Inches(6.7), Inches(2.05), Inches(6.1), Inches(4.5), size=15, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data & Baseline
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LGREY)
header_bar(s, "Data, Label, and Baseline", "What was measured and where we started")
slide_num(s, 3)

# Stats boxes
for i, (label, val, sub) in enumerate([
    ("786,420", "rows", "Chicago Traffic Tracker, Sept 2023"),
    ("3.49 : 1", "class ratio", "77.7% not-congested vs 22.3% congested"),
    ("~23 mph", "threshold", "30th percentile of training speeds → 'congested'"),
    ("30 min", "forecast horizon", "predict 3 time steps ahead per segment"),
]):
    x = Inches(0.3 + i * 3.2)
    rect(s, x, Inches(1.4), Inches(3.0), Inches(1.7), fill=NAVY)
    txbox(s, label, x, Inches(1.5), Inches(3.0), Inches(0.7),
          size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, val, x, Inches(2.15), Inches(3.0), Inches(0.4),
          size=14, bold=False, color=LBLUE, align=PP_ALIGN.CENTER)
    txbox(s, sub, x, Inches(2.55), Inches(3.0), Inches(0.5),
          size=11, color=LGREY, align=PP_ALIGN.CENTER)

# Baseline table
txbox(s, "Baseline: Logistic Regression (lags 1–3 + time features)",
      Inches(0.3), Inches(3.3), Inches(12.7), Inches(0.4),
      size=15, bold=True, color=NAVY)
mini_table(s, [
    ["", "F1", "Precision", "Recall", "Notes"],
    ["exp_001  Logistic Regression baseline", "0.5598", "0.7549", "0.4448", "Starting point"],
    ["exp_030  Final LightGBM (validation)", "0.6806", "0.6352", "0.7329", "+21.6% F1 gain"],
    ["exp_030  Final LightGBM (TEST SET)", "0.6183", "0.5639", "0.6843", "Held-out, once"],
], [Inches(4.6), Inches(0.85), Inches(1.05), Inches(0.95), Inches(2.5)],
   Inches(0.3), Inches(3.75), row_h=Inches(0.55), text_size=13)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Loop Design + Experiment Trace
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LGREY)
header_bar(s, "Loop Design and Experiment Trace",
           "How the AutoResearch loop was structured and what it explored")
slide_num(s, 4)

# Loop flow diagram (text-based boxes)
for i, (step, color) in enumerate([
    ("Read\nmodel.py", NAVY),
    ("Propose\none change", BLUE),
    ("Edit &\nrun", BLUE),
    ("F1\nimproved?", RGBColor(0x85,0x64,0x04)),
    ("Commit\n✓ keep", GREEN),
    ("Revert +\nlog discard", RED),
]):
    x = Inches(0.3 + i * 2.1)
    rect(s, x, Inches(1.4), Inches(1.85), Inches(1.2), fill=color)
    txbox(s, step, x, Inches(1.48), Inches(1.85), Inches(1.1),
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txbox(s, "→", Inches(0.3 + i*2.1 + 1.9), Inches(1.82), Inches(0.2), Inches(0.5),
              size=18, color=NAVY, align=PP_ALIGN.CENTER)

# Phase summary table
txbox(s, "30 experiments across 7 phases:", Inches(0.3), Inches(2.85), Inches(12.7), Inches(0.4),
      size=15, bold=True, color=NAVY)
mini_table(s, [
    ["Phase", "Experiments", "Peak F1", "Gain", "Verdict"],
    ["Model selection (LR → RF + balance + lags)", "exp_001–007", "0.6566", "+0.097", "RF wins"],
    ["RF hyperparameter grid (depth, trees, ratio)", "exp_008–013", "0.6585", "+0.002", "Ceiling hit"],
    ["Threshold sweep (0.35–0.50)", "exp_014–018", "0.6585", "+0.000", "Can't break it"],
    ["Segment-relative features added", "exp_019–021", "0.6727", "+0.014", "CEILING BROKEN"],
    ["Model comparison (RF/HGB/XGB/LightGBM)", "exp_022–025", "0.6780", "+0.005", "LightGBM wins"],
    ["LightGBM HP tuning (leaves, subsample)", "exp_026–030", "0.6806", "+0.003", "Final locked"],
], [Inches(4.0), Inches(1.8), Inches(0.85), Inches(0.75), Inches(3.6)],
   Inches(0.3), Inches(3.3), row_h=Inches(0.5), text_size=12)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Ceiling Discovery (key finding)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LGREY)
header_bar(s, "The Key Finding: A Labeling-Driven Ceiling",
           "Why hyperparameter tuning hit a wall — and how we fixed it")
slide_num(s, 5)

# Problem box
rect(s, Inches(0.3), Inches(1.4), Inches(6.1), Inches(5.5), fill=LRED)
txbox(s, "The Problem", Inches(0.5), Inches(1.5), Inches(5.7), Inches(0.45),
      size=17, bold=True, color=RED)
prob_items = [
    "Global threshold: speed < 23 mph → 'congested'",
    "But different roads have very different normal speeds",
    "Highway at 22 mph = normal rush-hour pace, labeled congested",
    "Side street at 20 mph = unusually fast, labeled not congested",
    "→ ~16.3% of training rows are mislabeled",
    "→ No model can learn from systematically wrong labels",
    "→ Hard ceiling at F1 ≈ 0.658 regardless of model or tuning",
]
bullet_box(s, prob_items, Inches(0.5), Inches(2.05), Inches(5.8), Inches(4.5),
           size=15, color=RED, gap=2)

# Fix box
rect(s, Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.5), fill=LGREEN)
txbox(s, "The Fix (no leakage)", Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.45),
      size=17, bold=True, color=GREEN)
fix_items = [
    "Compute per-segment stats from training only:",
    "  segment_mean_speed = mean(SPEED) per segment",
    "  segment_std_speed  = std(SPEED)  per segment",
    "  speed_zscore = (SPEED − mean) / std",
    "  speed_vs_seg_mean = SPEED / mean",
    "Merge onto both splits by SEGMENT_ID",
    "Now model asks: 'slow for THIS road?'",
    "F1: 0.6585 → 0.6727 in one step (+0.014)",
]
bullet_box(s, fix_items, Inches(6.9), Inches(2.05), Inches(5.8), Inches(4.5),
           size=15, color=GREEN, gap=2)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Final Result + Stability
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LGREY)
header_bar(s, "Final Result and Stability Check",
           "exp_030: LightGBM tuned — test set opened exactly once")
slide_num(s, 6)

# Big numbers
for i, (label, val, sub, col) in enumerate([
    ("Validation F1", "0.6806", "model selection metric", NAVY),
    ("TEST SET F1", "0.6183", "held-out, opened once", GREEN),
    ("vs Baseline", "+21.6%", "F1 gain, validation", BLUE),
    ("Stability std", "±0.005", "across 5 random seeds", DGREY),
]):
    x = Inches(0.3 + i * 3.2)
    rect(s, x, Inches(1.4), Inches(3.0), Inches(1.9), fill=col)
    txbox(s, label, x, Inches(1.48), Inches(3.0), Inches(0.45),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, val, x, Inches(1.9), Inches(3.0), Inches(0.75),
          size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, sub, x, Inches(2.65), Inches(3.0), Inches(0.5),
          size=11, color=LGREY, align=PP_ALIGN.CENTER)

# Final model config
rect(s, Inches(0.3), Inches(3.55), Inches(5.9), Inches(3.5), fill=NAVY)
txbox(s, "Final Model (exp_030)", Inches(0.5), Inches(3.65), Inches(5.5), Inches(0.4),
      size=14, bold=True, color=WHITE)
config = [
    "LightGBM  n_estimators=500, max_depth=6, lr=0.05",
    "num_leaves=31, min_child_samples=10",
    "colsample_bytree=0.8, subsample=0.8",
    "Undersample majority 2:1, threshold=0.40",
    "17 features (all segment-relative included)",
]
bullet_box(s, config, Inches(0.5), Inches(4.15), Inches(5.6), Inches(2.7),
           size=13, color=WHITE, gap=1)

# Stability table
txbox(s, "Stability across 5 seeds:", Inches(6.5), Inches(3.55), Inches(6.5), Inches(0.4),
      size=14, bold=True, color=NAVY)
mini_table(s, [
    ["Seed", "F1", "Precision", "Recall"],
    ["42",  "0.671", "0.595", "0.769"],
    ["0",   "0.668", "0.600", "0.755"],
    ["123", "0.674", "0.600", "0.768"],
    ["7",   "0.682", "0.636", "0.736"],
    ["99",  "0.674", "0.621", "0.737"],
    ["Mean ± Std", "0.674 ± 0.005", "", ""],
], [Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1)],
   Inches(6.5), Inches(4.0), row_h=Inches(0.42), text_size=13)

txbox(s, "Val→Test gap (−0.062): test period has 24.4% congestion vs 29.7% in validation\n"
         "— distribution shift in the later time window. No additional tuning was done.",
      Inches(6.5), Inches(6.95), Inches(6.5), Inches(0.5),
      size=11, color=DGREY, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Worked vs Failed
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LGREY)
header_bar(s, "What Worked vs. What Failed",
           "Honest breakdown of all 30 experiments")
slide_num(s, 7)

# Worked
rect(s, Inches(0.3), Inches(1.4), Inches(6.1), Inches(5.6), fill=LGREEN)
txbox(s, "✓  Worked", Inches(0.5), Inches(1.5), Inches(5.7), Inches(0.45),
      size=17, bold=True, color=GREEN)
worked = [
    "**Segment-relative features** → +0.014 F1 (single biggest gain)",
    "**Class imbalance: undersample 2:1** → cleaner precision than class_weight",
    "**LightGBM over RF** → consistent +0.003 on same setup",
    "**Subsampling regularization** → colsample+subsample=0.8 → +0.003",
    "**Frozen evaluator** → all 30 comparisons are valid and comparable",
    "**Git discipline** → zero lost experiments, zero accidental overrides",
]
bullet_box(s, worked, Inches(0.5), Inches(2.05), Inches(5.8), Inches(4.8),
           size=15, color=GREEN, gap=3)

# Failed
rect(s, Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.6), fill=LRED)
txbox(s, "✗  Did Not Work", Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.45),
      size=17, bold=True, color=RED)
failed = [
    "**HistGradientBoosting** → failed both attempts (exp_005, exp_022)",
    "**Threshold below 0.40** → raises recall, always drops F1",
    "**RF hyperparameter tuning** → 6 experiments, total gain < 0.003",
    "**XGBoost** → competitive but LightGBM wins every comparison",
    "**Fixing the label definition** → requires unfreezing run.py, out of scope",
    "**Single-seed evaluation** → kept lucky runs without knowing it",
]
bullet_box(s, failed, Inches(6.9), Inches(2.05), Inches(5.8), Inches(4.8),
           size=15, color=RED, gap=3)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Reflection + Limits
# ════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, NAVY)
header_bar(s, "Reflection: Limits of the AutoResearch Loop", "")
slide_num(s, 8)

txbox(s, "What did I learn about doing research with AI agents?",
      Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.45),
      size=18, bold=True, color=LBLUE, align=PP_ALIGN.LEFT)

# Four boxes
boxes = [
    ("Agent excelled at", [
        "Systematic bookkeeping — 30 experiments, 0 errors",
        "Execution speed — hours not weeks",
        "Identifying the labeling problem (measured 16.3% disagreement)",
        "Clean model comparison (held all variables constant)",
    ], BLUE, WHITE),
    ("Agent struggled with", [
        "No lookahead — committed to RF for 20 exps before trying LightGBM",
        "No repeat-on-change — HGB should have been retried with new features",
        "Single-seed decisions — lucky runs weren't flagged",
        "Can't diagnose *why* a ceiling exists — only measure it",
    ], RGBColor(0x6c,0x1f,0x1f), LRED),
    ("Human judgment was irreplaceable for", [
        "Framing 'why is the ceiling there?' — not just measuring it",
        "Deciding when diminishing returns justify stopping",
        "Keeping run.py frozen — protects all comparisons retroactively",
        "Choosing which result is the 'story' worth telling",
    ], RGBColor(0x1d,0x5c,0x2e), LGREEN),
    ("Fundamental limits", [
        "Congestion label is still globally defined — partially wrong labels remain",
        "Val→Test gap (−0.062): later time period has different congestion rate",
        "Single dataset, single month — generalizability unknown",
        "Greedy loop misses global optima by design",
    ], RGBColor(0x5c,0x45,0x04), LAMBER),
]

for i, (title, items, fill, txt_col) in enumerate(boxes):
    x = Inches(0.3 + (i % 2) * 6.5)
    y = Inches(1.9 + (i // 2) * 2.52)
    rect(s, x, y, Inches(6.2), Inches(2.35), fill=fill)
    txbox(s, title, x + Inches(0.15), y + Inches(0.1), Inches(5.9), Inches(0.38),
          size=13, bold=True, color=txt_col)
    bullet_box(s, items, x + Inches(0.15), y + Inches(0.5), Inches(5.9), Inches(1.75),
               size=12, color=txt_col, gap=1)

txbox(s, "Key insight: the loop shifts the bottleneck from execution to diagnosis — "
         "asking the right question matters more than running more experiments.",
      Inches(0.5), Inches(6.97), Inches(12.3), Inches(0.42),
      size=13, bold=True, color=LBLUE, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Slides written to {OUTPUT}  ({len(prs.slides)} slides)")
